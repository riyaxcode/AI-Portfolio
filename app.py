import streamlit as st
import pickle
import numpy as np
import pandas as pd
import requests
import nltk
import re
import PyPDF2
import docx
import io
import time
from pptx import Presentation
from datetime import date
import datetime as dt

# --- IMPORTS FOR STOCK PREDICTION ---
import matplotlib.pyplot as plt
import yfinance as yf
from keras.models import load_model
from sklearn.preprocessing import MinMaxScaler

# --- IMPORTS FOR HOUSE PRICE PREDICTION ---
from sklearn.datasets import fetch_california_housing
from sklearn.model_selection import train_test_split
from xgboost import XGBRegressor
from sklearn import metrics

st.set_page_config(page_title="AI Portfolio Hub", layout="wide")

# NLTK Downloads for Resume Screening
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)

# ==========================================
# --- Custom CSS (For Cards, Buttons & Headings) ---
# ==========================================
st.markdown("""
    <style>
    /* 1. HEADING STYLES (Uppercase, Center, Century Gothic) */
    h1, h2, h3, h4, h5, h6 {
        text-transform: uppercase !important;
        font-family: "Century Gothic", sans-serif !important;
        text-align: center !important;
    }

    /* 2. MOVIE CARD STYLES */
    .movie-container {
        text-align: center;
        margin-bottom: 20px;
    }
    .movie-img {
        border-radius: 10px;
        transition: transform 0.3s ease, box-shadow 0.3s ease;
        width: 100%;
        cursor: pointer;
    }
    .movie-img:hover {
        transform: scale(1.05);
        box-shadow: 0 10px 20px rgba(0,0,0,0.3);
    }
    .movie-title {
        font-size: 14px;
        font-weight: 600;
        margin-bottom: 5px;
        height: 40px;
        overflow: hidden;
        text-overflow: ellipsis;
        display: -webkit-box;
        -webkit-line-clamp: 2;
        -webkit-box-orient: vertical;
    }

    /* 3. BUTTON STYLES (Red Color) */
    div.stButton > button {
        background-color: #FF0000;
        color: white;
        border: none;
        border-radius: 8px;
        padding: 10px 20px;
        font-weight: bold;
        transition: background-color 0.3s;
    }
    div.stButton > button:hover {
        background-color: #CC0000; 
        color: white;
        border: none;
    }
    div.stButton > button:active {
        background-color: #990000;
        color: white;
    }
    </style>
    """, unsafe_allow_html=True)


# ==========================================
# --- Cached Model Loading Functions ---
# ==========================================
@st.cache_resource
def load_crop_models():
    dtr = pickle.load(open('Crop Yield Prediction/dtr.pkl', 'rb'))
    preprocessor = pickle.load(open('Crop Yield Prediction/preprocessor.pkl', 'rb'))
    return dtr, preprocessor


@st.cache_resource
def load_movie_data():
    with open('Movie_recommendor_system/movie_dict.pkl', 'rb') as f:
        movies_dict = pickle.load(f)
        movies = pd.DataFrame(movies_dict)
    with open('Movie_recommendor_system/similarity.pkl', 'rb') as f:
        similarity = pickle.load(f)
    return movies, similarity


@st.cache_resource
def load_resume_models():
    clf = pickle.load(open('Resume Screening Model/model.pkl', 'rb'))
    tfidf = pickle.load(open('Resume Screening Model/tfidf.pkl', 'rb'))
    le = pickle.load(open('Resume Screening Model/encoder.pkl', 'rb'))
    return clf, tfidf, le


@st.cache_resource
def load_and_train_house_model():
    house_price_dataset = fetch_california_housing()
    house_price_dataframe = pd.DataFrame(house_price_dataset.data, columns=house_price_dataset.feature_names)
    house_price_dataframe['price'] = house_price_dataset.target

    X = house_price_dataframe.drop(['price'], axis=1)
    Y = house_price_dataframe['price']
    X_train, X_test, Y_train, Y_test = train_test_split(X, Y, test_size=0.2, random_state=2)

    model = XGBRegressor()
    model.fit(X_train, Y_train)

    test_predictions = model.predict(X_test)
    r2 = metrics.r2_score(Y_test, test_predictions)
    mae = metrics.mean_absolute_error(Y_test, test_predictions)

    return model, r2, mae


# ==========================================
# --- Core Logic Functions ---
# ==========================================

# 1. Movie Recommendation Logic
def fetch_poster(movie_id):
    try:
        url = f"https://api.themoviedb.org/3/movie/{movie_id}?api_key=429d8ef2a929d522ee1b9cb1043e6961&language=en-US"
        response = requests.get(url, timeout=5)
        response.raise_for_status()
        data = response.json()
        poster_path = data.get('poster_path')

        if poster_path:
            return "https://image.tmdb.org/t/p/w500" + poster_path
        else:
            return "https://via.placeholder.com/500x750?text=No+Image"
    except requests.exceptions.RequestException:
        return "https://via.placeholder.com/500x750?text=Error"


def recommend_movies(movie_title, movies_df, similarity_matrix):
    try:
        movie_index = movies_df[movies_df['title'] == movie_title].index[0]
        distances = similarity_matrix[movie_index]
        movies_list = sorted(list(enumerate(distances)), reverse=True, key=lambda x: x[1])[1:6]

        recommended_movie_titles = []
        recommended_movie_posters = []

        for i in movies_list:
            movie_id = movies_df.iloc[i[0]].movie_id
            recommended_movie_titles.append(movies_df.iloc[i[0]].title)
            recommended_movie_posters.append(fetch_poster(movie_id))

        return recommended_movie_titles, recommended_movie_posters
    except Exception as e:
        st.error(f"An error occurred in recommendation logic: {e}")
        return [], []


# 2. Resume Screening Logic
def cleanResume(txt):
    cleanText = re.sub('http\\S+\\s', ' ', txt)
    cleanText = re.sub('RT|cc', ' ', cleanText)
    cleanText = re.sub('#\\S+\\s', ' ', cleanText)
    cleanText = re.sub('@\\S+', ' ', cleanText)
    cleanText = re.sub('[%s]' % re.escape("""!"#$%&'()*+,-./:;<=>?@[\\]^_`{|}~"""), ' ', cleanText)
    cleanText = re.sub(r'[^\\x00-\\x7f]', ' ', cleanText)
    cleanText = re.sub('\\s+', ' ', cleanText)
    return cleanText


def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ''
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text
    return text


def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ''
    for para in doc.paragraphs:
        text += para.text + '\n'
    return text


def extract_text_from_txt(file):
    try:
        text = file.read().decode('utf-8')
    except:
        text = file.read().decode('latin-1')
    return text


def extract_text_from_pptx(file):
    ppt = Presentation(file)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def handle_file_upload(uploaded_file):
    file_extension = uploaded_file.name.split('.')[-1].lower()

    if file_extension == 'pdf':
        text = extract_text_from_pdf(uploaded_file)
    elif file_extension == 'docx':
        text = extract_text_from_docx(uploaded_file)
    elif file_extension == 'txt':
        text = extract_text_from_txt(uploaded_file)
    elif file_extension == 'pptx':
        text = extract_text_from_pptx(uploaded_file)
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, DOCX, TXT, or PPTX file.")
    return text


def predict_resume_category(resume_text, clf, tfidf):
    cleaned_text = cleanResume(resume_text)
    vectorized_text = tfidf.transform([cleaned_text])
    prediction = clf.predict(vectorized_text)[0]
    return prediction


# ==========================================
# --- Sidebar Navigation ---
# ==========================================
st.sidebar.title("🧠 AI Model Hub")
st.sidebar.write("Navigate through the projects:")

project_options = [
    "🏠 Home",
    "🌾 Crop Yield Prediction",
    "🏘️ House Price Prediction",
    "🎬 Movie Recommendation System",
    "📈 Stock Price Prediction",
    "📄 Resume Screening Model",
]

selected_project = st.sidebar.selectbox("Choose a Project", project_options)

# ==========================================
# --- App Routing & Pages ---
# ==========================================

if selected_project == "🏠 Home":
    st.markdown(
        "<h1 style='text-align: center; font-family: \"Century Gothic\", sans-serif;'>Welcome to My AI Portfolio</h1>",
        unsafe_allow_html=True)

    try:
        st.image("images/AI Portfolio hub.png", use_container_width=True)
    except Exception:
        st.warning("⚠️ Hub image not found. Make sure 'AI Portfolio hub.png' is inside the 'images' folder.")

    st.subheader("Central AI Command Center")
    st.write(
        "Welcome to the digital nervous system of my portfolio. Below you will find detailed insights into each model module. Use the sidebar on the left to activate and interact with any specific model.")
    st.divider()

    col1, col2 = st.columns(2)
    with col1:
        with st.expander("🌾 Crop Yield Prediction", expanded=True):
            st.write(
                "**About the Model:** Uses environmental data like rainfall and temperature to forecast agricultural output.")
            st.info("👈 Select 'Crop Yield Prediction' from the sidebar to activate the model.")
        with st.expander("🏘️ House Price Prediction"):
            st.write("**About the Model:** Estimates property values based on location, size, and market factors.")
        with st.expander("🎬 Movie Recommendation System"):
            st.write(
                "**About the Model:** An AI-powered system suggesting personalized movies based on user preferences.")

    with col2:
        with st.expander("📈 Stock Price Prediction"):
            st.write("**About the Model:** Time-series financial forecasting for stock trends using deep learning.")
        with st.expander("📄 Resume Screening Model"):
            st.write("**About the Model:** Automates HR processes by predicting job categories from uploaded resumes.")


# ==========================================
# 🌾 CROP YIELD PREDICTION SECTION
# ==========================================
elif selected_project == "🌾 Crop Yield Prediction":
    st.title("🌾 Crop Yield Prediction")

    try:
        st.image("images/crop.png", use_container_width=True)
    except Exception:
        st.warning("⚠️ 'crop.png' not found. Please ensure it is saved in the 'images' folder.")

    st.divider()
    st.markdown(
        "<p style='text-align: center;'>Enter the details of your farm and environment to predict the expected crop yield.</p>",
        unsafe_allow_html=True)
    st.divider()

    try:
        dtr, preprocessor = load_crop_models()

        col1, col2 = st.columns(2)

        with col1:
            year = st.number_input("Year", value=1990)
            rainfall = st.number_input("Average Rainfall (mm/year)", min_value=0.0, value=1485.0)
            pesticides = st.number_input("Pesticides Used (tonnes)", min_value=0.0, value=121.0)

        with col2:
            temp = st.number_input("Average Temperature (°C)", min_value=-10.0, max_value=60.0, value=16.37)
            area_list = ["Algeria", "Angola", "Argentina", "Armenia", "Australia", "Austria", "Azerbaijan", "Bahamas",
                         "Bahrain", "Bangladesh", "Belarus", "Belgium", "Botswana", "Brazil", "Bulgaria",
                         "Burkina Faso", "Burundi", "Cameroon", "Canada", "Central African Republic", "Chile",
                         "Colombia", "Croatia", "Denmark", "Dominican Republic", "Ecuador", "Egypt", "El Salvador",
                         "Eritrea", "Estonia", "Finland", "France", "Germany", "Ghana", "Greece", "Guatemala", "Guinea",
                         "Guyana", "Haiti", "Honduras", "Hungary", "India", "Indonesia", "Iraq", "Ireland", "Italy",
                         "Jamaica", "Japan", "Kazakhstan", "Kenya", "Latvia", "Lebanon", "Lesotho", "Libya",
                         "Lithuania", "Madagascar", "Malawi", "Malaysia", "Mali", "Mauritania", "Mauritius", "Mexico",
                         "Zambia", "Zimbabwe"]
            area = st.selectbox("Area", area_list)
            crop_list = ["Maize", "Plantains and others", "Potatoes", "Rice, paddy", "Sorghum", "Soybeans",
                         "Sweet potatoes", "Wheat", "Yams"]
            item = st.selectbox("Crop Type", crop_list)

        st.divider()

        if st.button("Predict Crop Yield 🚀", use_container_width=True):
            with st.spinner("Calculating predictions..."):
                try:
                    input_data = pd.DataFrame([[year, rainfall, pesticides, temp, area, item]],
                                              columns=['Year', 'average_rain_fall_mm_per_year', 'pesticides_tonnes',
                                                       'avg_temp', 'Area', 'Item'])
                    transformed_features = preprocessor.transform(input_data)
                    prediction = dtr.predict(transformed_features)
                    st.success(f"### 🎉 Predicted Yield: {prediction[0]:.2f} hg/ha")
                except Exception as e:
                    st.error(f"❌ An error occurred during prediction: {e}")

    except FileNotFoundError as e:
        st.error(
            f"❌ Error loading models. Please ensure 'dtr.pkl' and 'preprocessor.pkl' are inside the 'Crop Yield Prediction' folder. Details: {e}")


# ==========================================
# 🎬 MOVIE RECOMMENDATION SECTION
# ==========================================
elif selected_project == "🎬 Movie Recommendation System":
    st.title("🎬 Movie Recommendation System")

    try:
        st.image("images/movie.png", use_container_width=True)
    except Exception:
        st.warning("⚠️ Conceptual system banner image ('movie.png') not found. Check the 'images' folder.")

    st.divider()
    st.markdown(
        "<p style='text-align: center;'>Our AI analyzes content similarity to recommend 5 movies you will love based on your selection.</p>",
        unsafe_allow_html=True)
    st.divider()

    try:
        with st.spinner("Loading Movie Database and Similarity Matrix..."):
            movies_df, similarity_matrix = load_movie_data()
        st.success("✅ Movie Recommendation Matrix Loaded Successfully!")

        movie_title_list = movies_df['title'].values

        selected_movie = st.selectbox(
            "🎥 Choose a movie you like from the dropdown:",
            movie_title_list,
            index=None,
            placeholder="Search or Select a Movie..."
        )

        if st.button("Recommend", use_container_width=True):
            if selected_movie:
                with st.spinner(f"Fetching recommendations for **{selected_movie}**..."):
                    names, posters = recommend_movies(selected_movie, movies_df, similarity_matrix)

                if names and posters:
                    st.subheader(f"If you liked **{selected_movie}**, you might enjoy these:")
                    cols = st.columns(5)
                    for col, name, poster in zip(cols, names, posters):
                        with col:
                            st.markdown(f"""
                                <div class="movie-container">
                                    <div class="movie-title">{name}</div>
                                    <img src="{poster}" class="movie-img">
                                </div>
                            """, unsafe_allow_html=True)
            else:
                st.warning("Please select a movie first.")

    except FileNotFoundError as e:
        st.error(f"❌ Error loading movie models. Detailed Error: {e}")
    except Exception as e:
        st.error(f"❌ An unexpected error occurred: {e}")

# 📄 RESUME SCREENING SECTION
# ==========================================
elif selected_project == "📄 Resume Screening Model":
    st.title("📄 Resume Screening Model")

    try:
        st.image("images/resumee.png", use_container_width=True)
    except Exception:
        st.warning("⚠️ 'resumee.png' not found. Please ensure it is saved in the 'images' folder.")

    st.divider()
    st.markdown(
        "<p style='text-align: center;'>Upload a resume in PDF, TXT, DOCX, or PPTX format to predict the job category.</p>",
        unsafe_allow_html=True)
    st.divider()

    try:
        # ⚠️ UPDATION 1: Yahan 'le' (LabelEncoder) ko bhi load karna hoga.
        # Ensure that your `load_resume_models()` function returns 3 things: clf, tfidf, le
        clf, tfidf, le = load_resume_models()

        uploaded_file = st.file_uploader("Upload Resume", type=["pdf", "docx", "txt", "pptx"])

        if uploaded_file is not None:
            try:
                resume_text = handle_file_upload(uploaded_file)
                st.success("File uploaded and read successfully!")

                with st.expander("View Extracted Text"):
                    st.write(resume_text)

                if st.button("Predict Category", use_container_width=True):
                    if resume_text.strip():
                        st.subheader("Prediction Results")
                        with st.spinner("Analyzing Resume..."):
                            # ⚠️ UPDATION 2: Yeh function ab ek number predict karega (e.g., 23)
                            prediction_number = predict_resume_category(resume_text, clf, tfidf)

                            # ⚠️ UPDATION 3: Number ko wapas Job Role String mein convert karna (Decoding)
                            # Agar output array form mein hai toh usko handle karne ke liye list [ ] use kiya hai
                            import numpy as np

                            if isinstance(prediction_number, (np.ndarray, list)):
                                category_name = le.inverse_transform(prediction_number)[0]
                            else:
                                category_name = le.inverse_transform([prediction_number])[0]

                        # Ab number ki jagah actual category name print hoga
                        st.success(f"### Predicted Job Category: **{category_name}**")
                    else:
                        st.warning("No readable text found in the document.")

            except Exception as e:
                st.error(f"An error occurred while reading or processing the file: {e}")

    except FileNotFoundError as e:
        st.error(
            f"❌ Error loading resume models. Make sure 'clf.pkl', 'tfidf.pkl', and 'encoder.pkl' are in the 'Resume Screening Model' folder. Details: {e}")

# ==========================================
# 📈 STOCK PRICE PREDICTION SECTION (UPDATED)
# ==========================================
elif selected_project == "📈 Stock Price Prediction":

    st.title("📈 Stock Price Prediction App")

    try:
        st.image("images/stock.png", use_container_width=True)
    except Exception:
        st.warning("⚠️ 'stock.png' not found. Please ensure it is saved in the 'images' folder.")

    st.divider()
    popular_stocks = [
        'POWERGRID.NS', 'RELIANCE.NS', 'TCS.NS', 'HDFCBANK.NS',
        'INFY.NS', 'ICICIBANK.NS', 'TATAMOTORS.NS', 'SBIN.NS',
        'BHARTIARTL.NS', 'ITC.NS', 'BAJFINANCE.NS', 'ASIANPAINT.NS',
        'AAPL', 'GOOGL', 'MSFT', 'TSLA', 'AMZN'
    ]

    selected_stock = st.selectbox(
        'Search or Select Stock Ticker',
        popular_stocks + ['Other']
    )

    if selected_stock == 'Other':
        user_input = st.text_input('Enter Custom Stock Ticker (e.g., WIPRO.NS)', 'WIPRO.NS')
    else:
        user_input = selected_stock
    # ----------------------------------------------------------------

    # Download the data
    start = dt.datetime(2000, 1, 1)
    end = dt.datetime(2024, 11, 1)

    st.subheader(f"Fetching Data for {user_input}...")
    try:
        df = yf.download(user_input, start=start, end=end)
        if df.empty:
            st.error("No data found for the given ticker. Please try a valid symbol.")
            st.stop()
    except Exception as e:
        st.error(f"Error fetching data: {e}")
        st.stop()

    st.subheader('Data from 2000 - 2024')
    st.write(df.describe())

    # Plotting the Closing Price
    st.subheader('Closing Price vs Time chart')
    fig = plt.figure(figsize=(12, 6))
    plt.plot(df['Close'], label='Closing Price', linewidth=1)
    plt.legend()
    st.pyplot(fig)

    # Plotting with 100 Days Moving Average
    st.subheader('Closing Price vs Time chart with 100MA')
    ma100 = df['Close'].rolling(100).mean()
    fig = plt.figure(figsize=(12, 6))
    plt.plot(df['Close'], label='Closing Price', linewidth=1)
    plt.plot(ma100, label='100 Days MA', linewidth=1.5)
    plt.legend()
    st.pyplot(fig)

    # Plotting with 100 Days & 200 Days Moving Averages
    st.subheader('Closing Price vs Time chart with 100MA & 200MA')
    ma200 = df['Close'].rolling(200).mean()
    fig = plt.figure(figsize=(12, 6))
    plt.plot(df['Close'], label='Closing Price', linewidth=1)
    plt.plot(ma100, label='100 Days MA', linewidth=1.5)
    plt.plot(ma200, label='200 Days MA', linewidth=1.5)
    plt.legend()
    st.pyplot(fig)

    # Data Splitting exactly like the notebook (70% Training, 30% Testing)
    data_training = pd.DataFrame(df['Close'][0:int(len(df) * 0.70)])
    data_testing = pd.DataFrame(df['Close'][int(len(df) * 0.70): int(len(df))])

    # Load Model
    st.subheader('Model Predictions')
    try:
        model = load_model('Stock Price Prediction/keras_model.keras')
    except Exception as e:
        st.error(f"Model file 'Stock Price Prediction/keras_model.keras' not found! Please ensure it is saved in the same directory.")
        st.stop()

    # Prepare testing data
    # We append the past 100 days of the training data to the testing data to predict the first test value
    past_100_days = data_training.tail(100)
    final_df = pd.concat([past_100_days, data_testing], ignore_index=True)

    scaler = MinMaxScaler(feature_range=(0, 1))
    input_data = scaler.fit_transform(final_df)

    x_test = []
    y_test = []

    for i in range(100, input_data.shape[0]):
        x_test.append(input_data[i - 100:i])
        y_test.append(input_data[i, 0])

    x_test, y_test = np.array(x_test), np.array(y_test)

    # Making Predictions
    y_predicted = model.predict(x_test)

    # Inverse transform to get actual values
    # We divide by the scaler scale factor to revert the normalization
    scale_factor = 1 / scaler.scale_[0]
    y_predicted = y_predicted * scale_factor
    y_test = y_test * scale_factor

    # Final Graph: Original vs Predicted Price
    st.subheader('Predictions vs Original')
    fig2 = plt.figure(figsize=(12, 6))
    plt.plot(y_test, label='Original Price', color='blue', linewidth=1.5)
    plt.plot(y_predicted, label='Predicted Price', color='red', linewidth=1.5)
    plt.xlabel('Time')
    plt.ylabel('Price')
    plt.legend()
    st.pyplot(fig2)

# ==========================================
# 🏘️ HOUSE PRICE PREDICTION SECTION
# ==========================================
elif selected_project == "🏘️ House Price Prediction":
    st.title("🏡 California House Price Predictor")

    try:
        st.image("images/house.png", use_container_width=True)
    except Exception:
        st.warning("⚠️ 'house.png' not found. Please ensure it is saved in the 'images' folder.")

    st.divider()
    st.markdown(
        "<p style='text-align: center;'>Enter the features of the house below to get an estimated price based on the California Housing Dataset.</p>",
        unsafe_allow_html=True)
    st.divider()

    with st.spinner("Loading dataset and training XGBoost model..."):
        model, r2, mae = load_and_train_house_model()

    st.sidebar.divider()
    st.sidebar.header("Model Performance")
    st.sidebar.write("Evaluation metrics on test data:")
    st.sidebar.write(f"**R² Score:** {r2:.4f}")
    st.sidebar.write(f"**Mean Absolute Error:** {mae:.4f}")

    st.subheader("Input House Features")

    col1, col2 = st.columns(2)

    with col1:
        MedInc = st.number_input("Median Income (in $10k)", min_value=0.0, max_value=20.0, value=3.87)
        HouseAge = st.number_input("House Age (Years)", min_value=1, max_value=100, value=28)
        AveRooms = st.number_input("Average Rooms", min_value=1, max_value=20, value=5)
        AveBedrms = st.number_input("Average Bedrooms", min_value=0, max_value=10, value=1)

    with col2:
        Population = st.number_input("Population", min_value=10.0, max_value=50000.0, value=1425.0)
        AveOccup = st.number_input("Average Occupancy", min_value=1.0, max_value=20.0, value=3.0)
        Latitude = st.number_input("Latitude", min_value=32.0, max_value=42.0, value=35.6)
        Longitude = st.number_input("Longitude", min_value=-125.0, max_value=-114.0, value=-119.5)

    st.divider()

    if st.button("Predict Price 🚀", use_container_width=True):
        input_data = np.array([[MedInc, HouseAge, AveRooms, AveBedrms, Population, AveOccup, Latitude, Longitude]])

        with st.spinner("Predicting Price..."):
            prediction = model.predict(input_data)

        estimated_price = prediction[0] * 100000

        st.success(f"### 🏡 Estimated House Price: ${estimated_price:,.2f}")